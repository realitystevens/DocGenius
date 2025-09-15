"""
Enhanced file processing service with improved text extraction and validation.
"""

import os
import io
import uuid
import hashlib
import mimetypes
from datetime import datetime
from typing import Dict, Any, List, Optional, BinaryIO

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from flask import current_app

from app.utils.exceptions import FileProcessingException


class FileService:
    """Enhanced file processing service with comprehensive text extraction."""

    def __init__(self):
        """Initialize the file service."""
        self.redis_client = None
        self.supported_extensions = {'.pdf', '.txt', '.docx', '.pptx'}
        self.max_file_size = 16 * 1024 * 1024  # 16MB

    def _get_redis_client(self):
        """Get Redis client from current app."""
        if not self.redis_client:
            self.redis_client = current_app.redis_client
        return self.redis_client

    def _redis_available(self):
        """Check if Redis is available."""
        return hasattr(current_app, 'redis_available') and current_app.redis_available

    def process_file(self, file, user_id: str) -> Dict[str, Any]:
        """
        Process an uploaded file and extract its content.

        Args:
            file: Uploaded file object
            user_id: User's session ID

        Returns:
            Dictionary containing file metadata and extracted content
        """
        # Generate unique file ID
        file_id = str(uuid.uuid4())

        # Read file content
        file_content = file.read()
        file.seek(0)  # Reset file pointer

        # Get file metadata
        file_name = file.filename
        file_size = len(file_content)
        file_extension = os.path.splitext(file_name)[1].lower()

        # Validate file size
        if file_size > self.max_file_size:
            raise FileProcessingException(
                f"File size ({file_size / 1024 / 1024:.1f}MB) exceeds maximum allowed size (16MB)"
            )

        # Extract text based on file type
        extracted_text = self._extract_text(
            io.BytesIO(file_content), file_extension)

        if not extracted_text or not extracted_text.strip():
            raise FileProcessingException(
                "No text content could be extracted from the file")

        # Calculate file hash for deduplication
        file_hash = hashlib.md5(file_content).hexdigest()

        # Prepare file data
        file_data = {
            'file_id': file_id,
            'file_name': file_name,
            'file_size': file_size,
            'file_extension': file_extension,
            'file_hash': file_hash,
            'extracted_text': extracted_text,
            'word_count': len(extracted_text.split()),
            'upload_timestamp': datetime.utcnow().isoformat(),
            'mime_type': mimetypes.guess_type(file_name)[0]
        }

        # Store in Redis - Redis is required for the application
        redis_client = self._get_redis_client()

        # Store file metadata
        redis_client.hset(f"file:{file_id}", mapping={
            'user_id': user_id,
            'file_name': file_name,
            'file_size': file_size,
            'file_extension': file_extension,
            'file_hash': file_hash,
            'word_count': file_data['word_count'],
            'upload_timestamp': file_data['upload_timestamp'],
            'mime_type': file_data['mime_type']
        })

        # Store extracted text separately for better performance
        redis_client.set(f"file_content:{file_id}", extracted_text)

        # Add to user's file list
        redis_client.lpush(f"user_files:{user_id}", file_id)

        # Set expiration (optional - 30 days)
        redis_client.expire(f"file:{file_id}", 30 * 24 * 60 * 60)
        redis_client.expire(f"file_content:{file_id}", 30 * 24 * 60 * 60)

        return {
            'file_id': file_id,
            'file_name': file_name,
            'file_size': file_size,
            'text_preview': extracted_text[:500] + ('...' if len(extracted_text) > 500 else ''),
            'word_count': file_data['word_count']
        }

    def _extract_text(self, file_stream: BinaryIO, file_extension: str) -> str:
        """
        Extract text from file based on its extension.

        Args:
            file_stream: File content as BytesIO
            file_extension: File extension (.pdf, .txt, etc.)

        Returns:
            Extracted text content
        """
        try:
            if file_extension == '.pdf':
                return self._extract_pdf_text(file_stream)
            elif file_extension == '.txt':
                return self._extract_txt_text(file_stream)
            elif file_extension == '.docx':
                return self._extract_docx_text(file_stream)
            elif file_extension == '.pptx':
                return self._extract_pptx_text(file_stream)
            else:
                raise FileProcessingException(
                    f"Unsupported file type: {file_extension}")
        except Exception as e:
            raise FileProcessingException(
                f"Failed to extract text from {file_extension} file: {str(e)}")

    def _extract_pdf_text(self, file_stream: BinaryIO) -> str:
        """Extract text from PDF file using PyMuPDF."""
        try:
            pdf_document = fitz.open(stream=file_stream, filetype="pdf")
            text_content = []

            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                text = page.get_text()
                if text.strip():
                    text_content.append(
                        f"--- Page {page_num + 1} ---\\n{text}")

            pdf_document.close()
            return "\\n\\n".join(text_content)

        except Exception as e:
            raise FileProcessingException(
                f"PDF text extraction failed: {str(e)}")

    def _extract_txt_text(self, file_stream: BinaryIO) -> str:
        """Extract text from TXT file."""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']

            for encoding in encodings:
                try:
                    file_stream.seek(0)
                    content = file_stream.read().decode(encoding)
                    return content
                except UnicodeDecodeError:
                    continue

            raise FileProcessingException(
                "Unable to decode text file with any supported encoding")

        except Exception as e:
            raise FileProcessingException(
                f"TXT text extraction failed: {str(e)}")

    def _extract_docx_text(self, file_stream: BinaryIO) -> str:
        """Extract text from DOCX file using python-docx."""
        try:
            document = Document(file_stream)
            text_content = []

            # Extract paragraphs
            for paragraph in document.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)

            # Extract tables
            for table in document.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    table_text.append(" | ".join(row_text))

                if table_text:
                    text_content.append("\\n".join(table_text))

            return "\\n\\n".join(text_content)

        except Exception as e:
            raise FileProcessingException(
                f"DOCX text extraction failed: {str(e)}")

    def _extract_pptx_text(self, file_stream: BinaryIO) -> str:
        """Extract text from PPTX file using python-pptx."""
        try:
            presentation = Presentation(file_stream)
            text_content = []

            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                if len(slide_text) > 1:  # More than just the slide header
                    text_content.append("\\n".join(slide_text))

            return "\\n\\n".join(text_content)

        except Exception as e:
            raise FileProcessingException(
                f"PPTX text extraction failed: {str(e)}")

    def get_user_files(self, user_id: str) -> List[Dict[str, Any]]:
        """
        Get all files for a specific user.

        Args:
            user_id: User's session ID

        Returns:
            List of file metadata dictionaries
        """
        files = []

        try:
            redis_client = self._get_redis_client()
            # Get user's file IDs
            file_ids = redis_client.lrange(f"user_files:{user_id}", 0, -1)

            for file_id in file_ids:
                file_data = redis_client.hgetall(f"file:{file_id}")
                if file_data:
                    files.append({
                        'file_id': file_id,
                        'file_name': file_data.get('file_name', ''),
                        'file_size': int(file_data.get('file_size', 0)),
                        'file_extension': file_data.get('file_extension', ''),
                        'word_count': int(file_data.get('word_count', 0)),
                        'upload_timestamp': file_data.get('upload_timestamp', ''),
                        'mime_type': file_data.get('mime_type', '')
                    })

            return files

        except Exception as e:
            current_app.logger.error(f"Error retrieving user files: {str(e)}")
            return []

    def get_file_content(self, file_id: str, user_id: str) -> Optional[Dict[str, Any]]:
        """
        Get file content and metadata.

        Args:
            file_id: File ID
            user_id: User's session ID

        Returns:
            File data dictionary or None if not found
        """
        try:
            redis_client = self._get_redis_client()
            # Check if file exists and belongs to user
            file_data = redis_client.hgetall(f"file:{file_id}")

            if not file_data or file_data.get('user_id') != user_id:
                return None

            # Get extracted text
            extracted_text = redis_client.get(f"file_content:{file_id}")

            if not extracted_text:
                return None

            return {
                'file_id': file_id,
                'file_name': file_data.get('file_name', ''),
                'extracted_text': extracted_text,
                'word_count': int(file_data.get('word_count', 0)),
                'upload_timestamp': file_data.get('upload_timestamp', '')
            }

        except Exception as e:
            current_app.logger.error(
                f"Error retrieving file content: {str(e)}")
            return None

    def delete_file(self, file_id: str, user_id: str) -> bool:
        """
        Delete a file and its associated data.

        Args:
            file_id: File ID
            user_id: User's session ID

        Returns:
            True if deleted successfully, False otherwise
        """
        redis_client = self._get_redis_client()

        try:
            # Check if file exists and belongs to user
            file_data = redis_client.hgetall(f"file:{file_id}")

            if not file_data or file_data.get('user_id') != user_id:
                return False

            # Delete file data
            redis_client.delete(f"file:{file_id}")
            redis_client.delete(f"file_content:{file_id}")

            # Remove from user's file list
            redis_client.lrem(f"user_files:{user_id}", 0, file_id)

            return True

        except Exception as e:
            current_app.logger.error(f"Error deleting file: {str(e)}")
            return False

    def get_file_count(self, user_id: str) -> int:
        """Get the number of files for a user."""
        redis_client = self._get_redis_client()
        try:
            return redis_client.llen(f"user_files:{user_id}")
        except Exception:
            return 0
