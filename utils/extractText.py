import fitz
from docx import Document
from pptx import Presentation




def extractPDFText(file_stream):
    try:
        pdf_document = fitz.open(stream=file_stream, filetype="pdf")
        text = ""

        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            text += page.get_text()
        return text
    except ValueError as e:
        return f"Extract PDF text Value Error: {e}"
    except Exception as e:
        return f"Extract PDF text Error: {e}" 
    

def extractTXTText(file_stream):
    try:
        file_stream.seek(0)  # Ensure the stream is at the beginning
        text = file_stream.read().decode('utf-8')
        return text
    except UnicodeDecodeError as e:
        return f"Extract TXT text Decode Error: {e}"
    except Exception as e:
        return f"Extract TXT text Error: {e}"
    

def extractDOCXText(file_stream):
    try:
        file_stream.seek(0)  # Ensure the stream is at the beginning
        doc = Document(file_stream)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text
    except Exception as e:
        return f"Extract DOCX text Error: {e}"
    

def extractPPTXText(file_stream):
    try:
        file_stream.seek(0)  # Ensure the stream is at the beginning
        presentation = Presentation(file_stream)
        text = ""

        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        return f"Extract PPTX text Error: {e}"
    